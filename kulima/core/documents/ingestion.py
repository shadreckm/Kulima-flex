"""Document ingestion service for Kulima OS (Phase A).

This module provides a vertical-agnostic DocumentIngestionService that can
accept uploaded files, detect their type, extract raw text, and produce
Document, DocumentChunk, and DocumentSource objects.

Phase A deliberately limits scope to:
- Basic text extraction for common formats
- Creation of models defined in `models.py`
- No claim extraction, scoring, or engine integration
"""

from __future__ import annotations

import io
import logging
import mimetypes
import os
import uuid
from typing import Iterable, List, Tuple

import pandas as pd
from pydantic import BaseModel

from kulima.core.documents.models import (
    Document,
    DocumentChunk,
    DocumentSource,
    DocumentType,
)
from kulima.models import SourceAttribution

_log = logging.getLogger(__name__)


class IngestedDocumentBundle(BaseModel):
    """Container for the outputs of a single ingestion call.

    This is a convenience wrapper for callers that want to treat
    documents, chunks, and sources as a single object.
    """

    document: Document
    chunks: List[DocumentChunk]
    sources: List[DocumentSource]


class DocumentIngestionService:
    """First-pass ingestion of uploaded documents (Phase A).

    This service is intentionally conservative: it aims to extract
    meaningful raw text while failing gracefully and logging errors when
    encountering unsupported or malformed files.
    """

    def infer_mime_type(self, filename: str) -> str:
        mime, _ = mimetypes.guess_type(filename)
        return mime or "application/octet-stream"

    def infer_doc_type(self, filename: str) -> DocumentType:
        name = filename.lower()
        if "deck" in name or name.endswith(".pptx"):
            return DocumentType.PITCH_DECK
        if "model" in name or name.endswith(".xlsx"):
            return DocumentType.FINANCIAL_MODEL
        if "minutes" in name:
            return DocumentType.BOARD_MINUTES
        if name.endswith(".csv") or name.endswith(".xlsx"):
            return DocumentType.DATA_EXPORT
        return DocumentType.GENERIC

    # ── Public API ───────────────────────────────────────────────────────

    def ingest_files(
        self,
        files: Iterable[object],
        *,
        run_id: int | None = None,
        uploaded_by: str | None = None,
        entities: list[str] | None = None,
    ) -> List[IngestedDocumentBundle]:
        """Ingest a collection of uploaded files.

        `files` are expected to behave like Streamlit UploadedFile objects:
        they must provide `.name` and `.read()` / `.getvalue()` to access
        bytes.  The service does not depend on Streamlit directly.
        """

        bundles: List[IngestedDocumentBundle] = []
        for f in files:
            try:
                bundle = self._ingest_single_file(
                    f,
                    run_id=run_id,
                    uploaded_by=uploaded_by,
                    entities=entities or [],
                )
                if bundle is not None:
                    bundles.append(bundle)
            except Exception as exc:  # pragma: no cover - defensive logging only
                _log.error(
                    "DocumentIngestionService: failed to ingest %s — %s: %s",
                    getattr(f, "name", "<unknown>"),
                    type(exc).__name__,
                    exc,
                )
        return bundles

    # ── Internal helpers ─────────────────────────────────────────────────

    def _ingest_single_file(
        self,
        file_obj: object,
        *,
        run_id: int | None,
        uploaded_by: str | None,
        entities: list[str],
    ) -> IngestedDocumentBundle | None:
        filename = getattr(file_obj, "name", "uploaded_file")
        mime_type = self.infer_mime_type(filename)
        doc_type = self.infer_doc_type(filename)
        doc_id = str(uuid.uuid4())

        raw_bytes = self._read_bytes(file_obj)
        text_chunks = self._extract_text_chunks(filename, mime_type, raw_bytes)
        if not text_chunks:
            _log.warning("DocumentIngestionService: no text extracted from %s", filename)

        document = Document(
            id=doc_id,
            filename=filename,
            mime_type=mime_type,
            doc_type=doc_type,
            uploaded_by=uploaded_by,
            entities=entities,
            metadata={"run_id": run_id},
        )

        chunks: List[DocumentChunk] = []
        for seq, (page_num, section, text) in enumerate(text_chunks):
            chunks.append(
                DocumentChunk(
                    id=str(uuid.uuid4()),
                    document_id=doc_id,
                    sequence=seq,
                    page_number=page_num,
                    section_heading=section,
                    text=text,
                )
            )

        # For Phase A we expose a single SourceAttribution per document,
        # using the filename as title and a synthetic URL.
        source = SourceAttribution(
            title=filename,
            url=f"document://{doc_id}",
            snippet=(chunks[0].text[:200] if chunks else ""),
            relevance=1.0,
            source_type="document",
            confidence_score=0.8,
        )
        doc_source = DocumentSource(document_id=doc_id, source=source)

        return IngestedDocumentBundle(document=document, chunks=chunks, sources=[doc_source])

    def _read_bytes(self, file_obj: object) -> bytes:
        # Streamlit's UploadedFile supports .getvalue(), but we also fall
        # back to .read() for generality.
        if hasattr(file_obj, "getvalue"):
            return file_obj.getvalue()
        if hasattr(file_obj, "read"):
            return file_obj.read()
        raise TypeError("Unsupported file object: missing read()/getvalue()")

    def _extract_text_chunks(
        self,
        filename: str,
        mime_type: str,
        raw_bytes: bytes,
    ) -> List[Tuple[int | None, str | None, str]]:
        """Extract a list of (page_number, section_heading, text) tuples.

        Phase A uses simple heuristics and standard libraries only.  It is
        intentionally best-effort and will log on failure rather than
        raising, returning an empty list when nothing can be extracted.
        """

        ext = os.path.splitext(filename)[1].lower()
        try:
            if ext in {".txt", ""} or mime_type.startswith("text/"):
                text = raw_bytes.decode("utf-8", errors="replace")
                return [(None, None, text)] if text.strip() else []

            if ext == ".pdf":
                try:
                    import PyPDF2  # type: ignore

                    reader = PyPDF2.PdfReader(io.BytesIO(raw_bytes))
                    chunks: list[Tuple[int | None, str | None, str]] = []
                    for i, page in enumerate(reader.pages):
                        page_text = page.extract_text() or ""
                        if page_text.strip():
                            chunks.append((i + 1, None, page_text))
                    return chunks
                except Exception as exc:  # pragma: no cover - optional dependency
                    _log.warning("PDF extraction failed for %s: %s", filename, exc)
                    return []

            if ext in {".docx"}:
                try:
                    import docx  # type: ignore

                    doc = docx.Document(io.BytesIO(raw_bytes))
                    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                    text = "\n".join(paragraphs)
                    return [(None, None, text)] if text.strip() else []
                except Exception as exc:  # pragma: no cover - optional dependency
                    _log.warning("DOCX extraction failed for %s: %s", filename, exc)
                    return []

            if ext in {".csv", ".tsv"}:
                try:
                    sep = "," if ext == ".csv" else "\t"
                    df = pd.read_csv(io.BytesIO(raw_bytes), sep=sep)
                    text = df.to_csv(index=False)
                    return [(None, None, text)] if text.strip() else []
                except Exception as exc:  # pragma: no cover
                    _log.warning("CSV/TSV extraction failed for %s: %s", filename, exc)
                    return []

            if ext in {".xlsx", ".xls"}:
                try:
                    xls = pd.ExcelFile(io.BytesIO(raw_bytes))
                    chunks: list[Tuple[int | None, str | None, str]] = []
                    for sheet_name in xls.sheet_names:
                        df = xls.parse(sheet_name)
                        text = df.to_csv(index=False)
                        if text.strip():
                            chunks.append((None, sheet_name, text))
                    return chunks
                except Exception as exc:  # pragma: no cover
                    _log.warning("Excel extraction failed for %s: %s", filename, exc)
                    return []

            if ext in {".pptx"}:
                try:
                    from pptx import Presentation  # type: ignore

                    prs = Presentation(io.BytesIO(raw_bytes))
                    chunks: list[Tuple[int | None, str | None, str]] = []
                    for i, slide in enumerate(prs.slides):
                        texts = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                texts.append(shape.text)
                        slide_text = "\n".join(t.strip() for t in texts if t.strip())
                        if slide_text:
                            chunks.append((i + 1, None, slide_text))
                    return chunks
                except Exception as exc:  # pragma: no cover
                    _log.warning("PPTX extraction failed for %s: %s", filename, exc)
                    return []

        except Exception as exc:  # pragma: no cover
            _log.warning(
                "DocumentIngestionService: generic extraction failure for %s: %s",
                filename,
                exc,
            )

        return []
